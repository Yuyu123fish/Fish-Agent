from __future__ import annotations

import io
import unittest

from pptx import Presentation

from fish_worker.parser.pptx_parser import PptxParser


def make_pptx_bytes() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "第一页标题"
    text_box = slide.shapes.add_textbox(left=0, top=0, width=1000000, height=1000000)
    text_box.text_frame.text = "正文内容"

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


class PptxParserTest(unittest.TestCase):
    def test_extracts_slide_text_with_page_number(self) -> None:
        elements = PptxParser().parse(make_pptx_bytes(), "deck.pptx")

        self.assertEqual([("第一页标题", 1), ("正文内容", 1)], [(e.text, e.page) for e in elements])
        self.assertTrue(all(e.element_type == "Text" for e in elements))


if __name__ == "__main__":
    unittest.main()
