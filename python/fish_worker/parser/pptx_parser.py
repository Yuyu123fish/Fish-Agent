"""PowerPoint .pptx parser based on python-pptx."""

from __future__ import annotations

import io

from fish_worker.parser.base import BaseParser, RawElement


class PptxParser(BaseParser):
    """Extract visible text frames from each slide."""

    def parse(self, content: bytes, filename: str, tmp_dir: str | None = None) -> list[RawElement]:
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(content))
        elements: list[RawElement] = []
        for slide_idx, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        elements.append(RawElement(text=text, page=slide_idx, element_type="Text"))
        return elements
